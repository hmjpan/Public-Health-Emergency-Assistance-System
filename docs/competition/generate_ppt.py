#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
卫盾Agent - GOAI初赛方案PPT (v3 修订版)
聚焦：7 Agent + 12 Skill | 双证据验证 | Guardrail | 证据链 | 法规溯源 | 188测试 | 双模式
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ========== 配色 ==========
C_DARK = RGBColor(0x1A, 0x3C, 0x6E)
C_MAIN = RGBColor(0x2B, 0x5C, 0x9A)
C_LIGHT = RGBColor(0x4A, 0x90, 0xD9)
C_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_RED = RGBColor(0xC6, 0x28, 0x28)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BG = RGBColor(0xF5, 0xF7, 0xFA)
C_LBLUE = RGBColor(0x90, 0xCA, 0xF9)
C_SOFT = RGBColor(0xBB, 0xDE, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL = 28

# ========== 工具函数 ==========
def bg(slide, color=C_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=18, color=C_GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = '微软雅黑'; p.alignment = align
    return tb

def mtxt(slide, l, t, w, h, lines, sz=15, color=C_GRAY):
    """多行文本"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = '微软雅黑'; p.space_after = Pt(6)
    return tb

def tbl(slide, l, t, w, h, data, hcolor=C_MAIN):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    table = ts.table
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = table.cell(i, j); c.text = str(val)
            p = c.text_frame.paragraphs[0]; p.font.size = Pt(13); p.font.name = '微软雅黑'
            if i == 0:
                p.font.color.rgb = C_WHITE; p.font.bold = True
                c.fill.solid(); c.fill.fore_color.rgb = hcolor
            else:
                p.font.color.rgb = C_GRAY
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE if i % 2 == 1 else RGBColor(0xF0, 0xF4, 0xF8)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    return ts

def pnum(slide, n):
    txt(slide, 12, 7.1, 1.2, 0.3, f"{n}/{TOTAL}", sz=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

def chapter(num, title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C_DARK)
    txt(s, 1, 2.5, 11, 1, f"第{num}章", sz=24, color=C_LBLUE)
    txt(s, 1, 3.3, 11, 1.2, title, sz=40, color=C_WHITE, bold=True)
    if sub: txt(s, 1, 4.5, 11, 0.8, sub, sz=18, color=C_LBLUE)
    return s

def header(slide, title):
    rect(slide, 0, 0, 13.333, 0.08, C_MAIN)
    txt(slide, 0.5, 0.3, 10, 0.6, title, sz=28, color=C_DARK, bold=True)

# ==================== Slide 1: 封面 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C_DARK)
rect(s, 0, 0, 13.333, 0.15, C_ORANGE)
txt(s, 1, 1.5, 11, 0.6, "GOAI · 新智基座 Agent Infra · 初赛方案", sz=16, color=C_LBLUE)
txt(s, 1, 2.5, 11, 1.5, "卫盾Agent", sz=56, color=C_WHITE, bold=True)
txt(s, 1, 4.0, 11, 1, "突发公共卫生事件多部门协同应急处置智能平台", sz=28, color=C_SOFT)
txt(s, 1, 5.2, 11, 0.5, "【团队名称】", sz=18, color=C_LBLUE)
txt(s, 1, 5.9, 11, 0.8,
    "辅助处置（真实事件） + 模拟演练（单人练习）\n系统固化流程 · 大模型扩展能力 · 每次可改进",
    sz=14, color=RGBColor(0x78, 0x90, 0xA0))

# ==================== Slide 2: 目录 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, 0, 0, 0.15, 7.5, C_MAIN)
txt(s, 0.8, 0.5, 5, 0.8, "目录", sz=36, color=C_DARK, bold=True)
toc = [
    ("01", "场景与价值", "核心痛点 · 量化价值 · 差异化（双证据/Guardrail/证据链）"),
    ("02", "方案总览", "双层架构 · 7 Agent + 12 Skill · MCP · 可观测"),
    ("03", "多Agent协同设计", "7个Agent · 双模式运行 · AgentTeams编排 · 安全边界"),
    ("04", "Skills工程体系", "12个Skill · Guardrail约束 · 证据链 · 法规溯源"),
    ("05", "工程落地与安全审计", "188测试全绿 · 7 MCP工具 · Trace/Metrics"),
    ("06", "开放/开源计划", "SKILL.md/WORKER.md/AGENTS.md · 复用 · 依赖"),
    ("07", "落地计划与进展", "Agent/Skill层已100%完成 · 里程碑 · 风险"),
    ("08", "团队介绍", ""),
]
for i, (num, title, desc) in enumerate(toc):
    y = 1.5 + i * 0.7
    txt(s, 1, y, 0.6, 0.5, num, sz=22, color=C_ORANGE, bold=True)
    txt(s, 1.8, y, 4, 0.5, title, sz=18, color=C_DARK, bold=True)
    txt(s, 5.5, y, 7, 0.5, desc, sz=13, color=C_GRAY)
pnum(s, 2)

# ==================== 第一章 ====================
chapter("一", "场景与价值", "卫生系统多部门协作 · 核心痛点 · 量化价值 · 创新差异化")
pnum(s, 3)

# ==================== Slide 4: 痛点 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "目标用户与核心痛点")
txt(s, 0.5, 0.9, 12, 0.5, "目标用户：卫生系统、疾控中心、卫生监督所、急救中心——多部门协同作战", sz=15, color=C_MAIN)

pain = [
    ["痛点", "表现", "后果"],
    ["乱", "人员思想不统一\n各组行动缺乏统一协调，各干各的", "资源冲突、任务重叠或遗漏\n现场混乱失控"],
    ["慢", "靠电话逐级通知、人工层层传达\n从接报到各组响应动辄数小时", "贻误最佳处置窗口\n事态扩大"],
    ["信息不一致", "各组各自掌握碎片信息\n上报数据口径不统一，汇总靠人工", "指挥层看不到全局\n决策依据失真"],
    ["能力参差", "部分人员从未处理过类似事件\n不知道怎么做，工作无效闲置", "执行不到位、流程出错\n关键时刻掉链子"],
    ["经验误判", "高度依赖个人经验\n水平差异导致形势误判、指挥错误", "级别定错→资源配置失当\n指挥错误→二次风险"],
    ["培训低效", "培训演练周期长、效果差\n训完就忘，即使多次演练\n发生事件仍无法统一高效部署", "演练≠实战能力\n事件来了还是手忙脚乱"],
]
tbl(s, 0.3, 1.5, 12.5, 5, pain)
pnum(s, 4)

# ==================== Slide 5: 场景示例 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "真实场景：学校食物中毒多部门协同")

txt(s, 0.5, 1.1, 12, 0.4, "场景：某区学校20余名学生呕吐腹泻 → 多部门如何从「乱」到「有序」？", sz=16, color=C_ORANGE, bold=True)

left_lines = [
    "❌ 传统方式（乱+慢+能力参差+信息不一致）：",
    "  · 电话逐个通知各组 → 2小时才通知完",
    "  · 部分人员没经验 → 不知道怎么做，工作闲置",
    "  · 指挥员凭经验判断 → 低估级别，指挥错误",
    "  · 各组上报格式不一 → 汇总靠人工拼凑",
    "  · 多次演练训完就忘 → 事件来了还是手忙脚乱",
]
mtxt(s, 0.5, 1.6, 5.5, 3, left_lines, sz=14, color=C_RED)

right_lines = [
    "✅ 卫盾Agent（有序+快+能力标准化+信息一致）：",
    "  · Sentinel双证据分类→5min完成识别",
    "  · GradeEvaluate研判定级→法规条文引用",
    "  · 规则引擎+LLM辅助→弥补经验不足",
    "  · Event State共享→各组看到同一张图",
    "  · Review Agent自动复盘→改进回流知识库",
]
mtxt(s, 6.5, 1.6, 6.5, 3, right_lines, sz=14, color=C_GREEN)

# 流程图
flow = [
    ("接报", C_LIGHT), ("识别分类", C_LIGHT), ("研判定级", C_MAIN),
    ("一键启动", C_MAIN), ("多组并行", C_ORANGE), ("统一采集", C_ORANGE),
    ("复盘沉淀", C_GREEN),
]
for i, (label, color) in enumerate(flow):
    x = 0.5 + i * 1.8
    rect(s, x, 5.0, 1.5, 0.6, color)
    txt(s, x, 5.08, 1.5, 0.45, label, sz=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        txt(s, x + 1.5, 5.05, 0.3, 0.5, "→", sz=18, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(s, 0.5, 5.8, 12, 0.6, "全链路：信息统一 · 并行协调 · 大模型辅助决策 · 事后自动复盘改进", sz=14, color=C_MAIN, bold=True)
pnum(s, 5)

# ==================== Slide 6: 量化价值 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "可量化价值")

val = [
    ["指标", "传统方式", "卫盾Agent", "提升"],
    ["多组协调时间", "2-3小时（电话逐个通知）", "≤5分钟（Agent并行调度）", "↓96%"],
    ["信息一致性", "各组数据口径不一", "Event State统一共享", "100%一致"],
    ["研判定级准确率", "依赖个人经验（差异大）", "规则引擎+大模型辅助", "标准化"],
    ["新手上岗能力", "没经验就不知道怎么干", "岗位SOP+知识库实时指引", "不依赖个人"],
    ["培训转化率", "训完就忘，实战还是不会", "系统内置流程，战中即练", "持续可用"],
    ["经验复用", "每次从零开始", "复盘→知识库→Skill迭代", "持续进化"],
]
tbl(s, 0.5, 1.2, 12, 3.5, val)

txt(s, 0.5, 5.2, 12, 0.8,
    "核心价值：不是替代人，而是让系统和大模型协助人做好每一次处置，\n并通过复盘机制让每一次都成为下一次的改进起点。",
    sz=16, color=C_MAIN, bold=True)
pnum(s, 6)

# ==================== Slide 7: 差异化 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "创新点与差异化")

diff = [
    ["维度", "现有做法", "卫盾Agent"],
    ["决策辅助", "纯人工 or 纯AI黑盒", "确定性规则引擎+LLM双轨·可解释可审计"],
    ["可靠性", "单一LLM判断", "双证据验证(LLM+规则引擎交叉校验)"],
    ["安全护栏", "无约束或硬编码", "Guardrail四类约束·每个Skill独立定义"],
    ["法规合规", "人工记忆法规", "legal_ref溯源→具体条文可追溯"],
    ["证据链", "黑盒无解释", "provenance/methodology/confidence/missing_evidence"],
    ["培训演练", "周期长效果差", "双模式运行：辅助处置+模拟演练"],
    ["事后复盘", "写个总结就完了", "ReviewReport+证据链+改进回流知识库"],
]
tbl(s, 0.5, 1.2, 12, 4.5, diff)

txt(s, 0.5, 6.0, 12, 0.6,
    "核心差异：每个决策可解释、可审计、可追溯到法规条文——处置→复盘→改进→闭环进化",
    sz=15, color=C_ORANGE, bold=True)
pnum(s, 7)

# ==================== 第二章 ====================
chapter("二", "方案总览", "双层架构：系统固化流程 + 大模型扩展能力")
pnum(s, 8)

# ==================== Slide 9: 架构 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "总体架构：系统层 + Agent/Skill层 + MCP + 可观测")

# 系统层
rect(s, 0.3, 1.0, 12.7, 0.5, C_MAIN)
txt(s, 0.5, 1.05, 12, 0.4, "第一层：系统层 —— 规则引擎(v2026.2) · 状态机 · 通知引擎 · 5类事件包 · 数据硬化",
    sz=16, color=C_WHITE, bold=True)

sys_layers = [
    (1.6, "人的管理", "11类角色RBAC · 多通道通知+强制反馈闭环 · 指令必达", RGBColor(0xE3, 0xF2, 0xFD)),
    (2.2, "物的管理", "人员出动跟踪 · 物资调拨闭环 · 车辆调度 · 医疗资源可视", RGBColor(0xBB, 0xDE, 0xFB)),
    (2.8, "流程管理", "三阶段状态机 · 完成标准校验 · 标准任务包 · 法规条文引用(legal_ref)", RGBColor(0x90, 0xCA, 0xF9)),
    (3.4, "规则引擎", "4张规则表v2026.2 · 分级/时限/编成/知识 · 纯函数可审计 · 75条回归测试", RGBColor(0x64, 0xB5, 0xF6)),
]
for y, title, desc, color in sys_layers:
    rect(s, 0.5, y, 12.3, 0.5, color)
    txt(s, 0.8, y + 0.02, 2, 0.25, title, sz=13, color=C_DARK, bold=True)
    txt(s, 2.8, y + 0.02, 9.8, 0.25, desc, sz=12, color=C_GRAY)

# Agent/Skill层
rect(s, 0.3, 4.0, 12.7, 0.5, C_ORANGE)
txt(s, 0.5, 4.05, 12, 0.4, "第二层：Agent/Skill层 —— 7个Agent + 12个Skill(v1.1.0) · Guardrail + 证据链 · 双证据验证",
    sz=16, color=C_WHITE, bold=True)

ai_layers = [
    (4.6, "规则引擎包装(7个)", "TypePack · GradeEval · SlaCheck · StaffAssign · PlanMatch · CriteriaCheck · NotifyDispatch", RGBColor(0xFF, 0xE0, 0xB2)),
    (5.2, "LLM增强(5个)", "EventClassify · CaseTrack · SOPGuidance · ReviewReport · DocDraft", RGBColor(0xFF, 0xCC, 0x80)),
    (5.8, "工具+治理", "7个MCP工具(contacts/materials/vehicles/medical/knowledge/notification/event-state) · Trace+Metrics可观测", RGBColor(0xFF, 0xB7, 0x4D)),
]
for y, title, desc, color in ai_layers:
    rect(s, 0.5, y, 12.3, 0.5, color)
    txt(s, 0.8, y + 0.02, 2.5, 0.25, title, sz=13, color=C_DARK, bold=True)
    txt(s, 3.3, y + 0.02, 9.3, 0.25, desc, sz=12, color=C_GRAY)

txt(s, 0.5, 6.5, 12, 0.5,
    "确定性规则引擎（可审计） + LLM增强（Guardrail约束） + MCP标准协议 + Trace/Metrics = 完整Agent Infra",
    sz=14, color=C_ORANGE, bold=True)
pnum(s, 9)

# ==================== 第三章 ====================
chapter("三", "多Agent协同设计", "7个Agent · 双模式运行 · AgentTeams编排 · 安全边界")
pnum(s, 10)

# ==================== Slide 11: 七个Agent ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "七个职能Agent — 模拟真实多部门协同体系")

agents = [
    ["Agent", "对应部门", "核心职责", "专属Skill", "解决的痛点"],
    ["🛰️ Sentinel", "信息科/值班室", "统一接报→事件分类", "EventClassify", "慢（自动分类识别）"],
    ["🎯 Commander", "指挥中心", "研判+定级+启动+推进", "GradeEvaluate", "经验误判（辅助定级）"],
    ["📦 Dispatch", "后勤/装备科", "人员+物资+车辆调度", "StaffAssign", "乱+无法统一部署"],
    ["🏕️ Field", "流调/消杀/管控组", "任务分发+SOP指引+采集", "SOPGuidance", "能力参差（新手按SOP做）"],
    ["🏥 Medical", "医疗救治组", "病例跟踪+密接+资源", "CaseTrack", "信息不一致（统一台账）"],
    ["📡 Communication", "宣教/信息科", "通知+反馈+发布", "NotifyDispatch", "慢+信息不一致"],
    ["📋 Review", "质控/复盘组", "复盘+整改+知识沉淀", "ReviewReport", "培训低效+经验沉淀"],
]
tbl(s, 0.2, 1.1, 12.9, 5.5, agents)
pnum(s, 11)

# ==================== Slide 12: 协同流程 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "端到端协同流程")

txt(s, 0.5, 1.1, 12, 0.4, "多Agent如何从「乱」到「有序」完成一次突发公共卫生事件处置？", sz=15, color=C_ORANGE, bold=True)

steps = [
    ("1.统一接报", "Sentinel接收上报\n自动识别事件类型"),
    ("2.辅助研判", "Commander调用规则引擎\n+大模型分析→定级建议"),
    ("3.一键启动", "Commander确认启动\n按类型匹配通知组"),
    ("4.并行调度", "Dispatch统一调度\n人员/物资/车辆并行"),
    ("5.统一处置", "Field管理多组任务\n表单统一采集回传"),
    ("6.信息同步", "Event State共享\n各组看到同一张图"),
    ("7.复盘改进", "Review自动提取改进\n回流知识库→Skill升级"),
]
for i, (title, desc) in enumerate(steps):
    x = 0.3 + i * 1.85
    color = C_LIGHT if i % 2 == 0 else C_MAIN
    rect(s, x, 1.7, 1.7, 0.5, color)
    txt(s, x, 1.75, 1.7, 0.4, title, sz=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, 2.35, 1.7, 1.2, desc, sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(s, 0.5, 3.8, 12, 0.4, "信息共享机制：Event State", sz=18, color=C_DARK, bold=True)
mtxt(s, 0.5, 4.3, 12, 1.5, [
    "所有Agent读写同一个Event State共享状态容器，包含：事件信息、决策记录、任务状态、资源状态、通知记录",
    "→ 解决「信息不一致」：各组不再各记各的，指挥层看到实时全局视图",
    "→ 解决「乱」：Agent统一编排，任务不重叠不遗漏",
    "→ 解决「慢」：多Agent并行执行，不再逐级等待",
], sz=14, color=C_GRAY)

txt(s, 0.5, 5.8, 12, 0.6,
    "异常处理：启动未达标→Dispatch催办 | 任务阻塞→Commander介入 | 通知超时→自动升级",
    sz=13, color=C_MAIN, bold=True)
pnum(s, 12)

# ==================== Slide 13: 安全边界 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "安全边界与决策辅助机制")

safe = [
    ["动作", "风险", "Agent辅助方式", "人工角色"],
    ["事件定级", "🔴 高", "规则引擎量化评估+大模型分析建议", "指挥员最终确认"],
    ["一键启动", "🔴 高", "自动匹配通知组+编成建议", "指挥员审批启动"],
    ["阶段推进", "🟡 中", "CriteriaCheck校验完成标准", "指挥员确认推进"],
    ["信息发布", "🔴 高", "大模型辅助起草发布稿", "审批后方可发布"],
    ["处置方案", "🟡 中", "知识库匹配SOP+大模型建议", "现场指挥确认执行"],
    ["通知下发", "🟢 低", "Agent自动执行+强制反馈", "全自动，留痕可查"],
    ["复盘沉淀", "🟡 中", "Review Agent自动提取改进", "人工审核后入库"],
]
tbl(s, 0.3, 1.1, 12.7, 4.5, safe)

txt(s, 0.5, 5.9, 12, 0.8,
    "设计原则：大模型和规则引擎辅助决策，但最终决策权始终在人\n"
    "高风险动作→Agent提供依据和建议，人工拍板 | 低风险动作→Agent自动执行，全程留痕",
    sz=14, color=C_MAIN, bold=True)
pnum(s, 13)

# ==================== 第四章 ====================
chapter("四", "Skills工程体系", "12个Skill · Guardrail约束 · 证据链 · 法规溯源 · SKILL.md标准规格")
pnum(s, 14)

# ==================== Slide 15: Agent-Skill映射 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "12个Skill × 7个Agent — 标准化规格(v1.1.0)")

txt(s, 0.5, 1.0, 6, 0.4, "每个Agent至少一个专属Skill，每个Skill含Guardrail + 证据链 + 法规引用", sz=15, color=C_ORANGE, bold=True)

dedicated = [
    ["Agent", "专属Skill", "Skill用途 · 法规依据"],
    ["🛰️ Sentinel", "EventClassify", "双证据分类(LLM+规则) · 《报告管理工作规范》"],
    ["🎯 Commander", "GradeEvaluate", "研判定级→风险等级 · 国办函〔2006〕39号 1.3"],
    ["📦 Dispatch", "StaffAssign", "编成建议(工作组×级别乘数) · 卫办应急发〔2010〕74号"],
    ["🏕️ Field", "SOPGuidance", "CDC 10步法+场景SOP · CDC现场调查10步法"],
    ["🏥 Medical", "CaseTrack", "病例分析+流行病学指标 · 《传染病防治法》"],
    ["📡 Communication", "NotifyDispatch", "多通道+强制反馈+超时升级 · 国务院令第376号"],
    ["📋 Review", "ReviewReport", "复盘+改进建议 · 《应急预案》第5.3条"],
]
tbl(s, 0.2, 1.5, 7.5, 4.5, dedicated)

txt(s, 8.0, 1.0, 5, 0.4, "共享Skill（跨Agent复用）", sz=15, color=C_MAIN, bold=True)
shared = [
    ["Skill", "用途", "使用Agent"],
    ["TypePackResolve", "事件类型配置加载", "Sentinel,Field,Medical"],
    ["SlaCheck", "法定时限校验", "Commander"],
    ["PlanMatch", "预案/SOP知识库匹配", "Commander,Field"],
    ["CriteriaCheck", "阶段完成标准校验", "Commander"],
    ["DocDraft", "文档起草(简报/终报)", "Communication,Review"],
]
tbl(s, 8.0, 1.5, 5.0, 3.5, shared)

txt(s, 8.0, 5.3, 5, 1.5,
    "v1.1.0升级：\n· Guardrail四类约束\n· evidence证据链\n· 法规条文引用\n· 缺失证据检测",
    sz=12, color=C_GRAY)
pnum(s, 15)

# ==================== Slide 16: 双模式 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "双模式运行：辅助处置 + 模拟演练")

txt(s, 0.5, 1.1, 12, 0.4, "同一套系统，两种运行模式——真实事件辅助处置 & 单人模拟演练", sz=15, color=C_ORANGE, bold=True)

# 模式一
rect(s, 0.3, 1.7, 6.2, 0.5, C_MAIN)
txt(s, 0.5, 1.75, 6, 0.4, "模式一：辅助处置（真实事件）", sz=16, color=C_WHITE, bold=True)
real_lines = [
    "· 各岗位真人操作，Agent辅助增强",
    "· Sentinel Agent自动接报分类",
    "· Commander Agent辅助研判定级",
    "· Dispatch Agent并行调度资源",
    "· Field Agent提供SOP指引",
    "· Review Agent自动复盘沉淀",
    "· 人做决策，Agent提供依据和建议",
]
mtxt(s, 0.5, 2.3, 6, 3, real_lines, sz=14, color=C_MAIN)

# 模式二
rect(s, 6.8, 1.7, 6.2, 0.5, C_GREEN)
txt(s, 7.0, 1.75, 6, 0.4, "模式二：模拟演练（培训练习）", sz=16, color=C_WHITE, bold=True)
drill_lines = [
    "· 一人扮演某角色（如指挥员）",
    "· 其余Agent充当固定角色自动执行",
    "· Sentinel充当信息员自动上报",
    "· Dispatch充当后勤自动调度",
    "· Field充当各组自动执行任务",
    "· 一个人即可走完全流程",
    "· 结束后Review Agent自动复盘",
]
mtxt(s, 7.0, 2.3, 6, 3, drill_lines, sz=14, color=C_GREEN)

# 核心说明
txt(s, 0.5, 5.2, 12, 0.4, "关键设计", sz=18, color=C_DARK, bold=True)
mode_data = [
    ["", "辅助处置模式", "模拟演练模式"],
    ["人员", "各岗位真人", "1人 + Agent充当固定角色"],
    ["Agent角色", "辅助者（提供建议和信息）", "扮演者（代替真人执行）"],
    ["适用场景", "实际突发事件处置", "个人培训、岗前练习、流程考核"],
    ["核心价值", "弥补个人经验不足", "解决多人协调难、演练频次低"],
]
tbl(s, 0.3, 5.6, 12.5, 2, mode_data)
pnum(s, 16)

# ==================== Slide 16b: 大模型辅助 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "三层能力体系 + Guardrail + 双证据验证")

txt(s, 0.5, 1.1, 12, 0.4, "确定性规则引擎 + LLM增强 + 双证据交叉验证 = 高可靠Agent决策", sz=16, color=C_ORANGE, bold=True)

dual = [
    ["层级", "方式", "Skill示例", "法规依据"],
    ["D0 规则层", "确定性规则引擎\n纯函数，同输入同输出", "GradeEvaluate\nSlaCheck\nStaffAssign", "国办函〔2006〕39号\n国务院令第376号\n卫办应急发〔2010〕74号"],
    ["D1 检索层", "知识库加权匹配\n确定性检索+评分排序", "PlanMatch\nTypePackResolve", "各类技术方案\n《报告管理工作规范》"],
    ["D2 LLM层", "LLM语义分析\n+规则引擎交叉校验", "EventClassify\nSOPGuidance\nCaseTrack", "双证据验证:\nLLM + 规则引擎\n置信度校准"],
]
tbl(s, 0.3, 1.7, 8.0, 3.5, dual)

# Guardrail box
rect(s, 8.6, 1.7, 4.4, 3.5, RGBColor(0xFD, 0xF0, 0xE0))
txt(s, 8.8, 1.8, 4, 0.3, "Guardrail 四类约束", sz=14, color=C_ORANGE, bold=True)
guard_lines = [
    "🔒 确定性约束",
    "  纯函数/同输入同输出",
    "📋 业务约束",
    "  就高不就低/白名单变量",
    "🛡️ 容错约束",
    "  规则未覆盖→转人工/不猜测",
    "⚡ 安全约束",
    "  结果仅建议/人工最终确认",
]
mtxt(s, 8.8, 2.2, 4, 3, guard_lines, sz=11, color=C_GRAY)

txt(s, 0.5, 5.5, 12, 1,
    "证据链：每个Skill输出含 provenance(来源) + methodology(方法) + confidence(置信度) + missing_evidence(缺失)\n"
    "→ 所有决策可解释、可审计、可追溯到具体法规条文",
    sz=14, color=C_MAIN, bold=True)
pnum(s, 17)

# ==================== Slide 17: 复盘改进机制 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "每次可改进可复盘的系统机制")

txt(s, 0.5, 1.1, 12, 0.4, "核心亮点：每次事件处置都是下一次改进的起点", sz=16, color=C_ORANGE, bold=True)

# 闭环图
loop = [
    (1.5, "事件处置", "7个Agent协同\n完成全流程", C_MAIN),
    (4.5, "自动复盘", "Review Agent提取\n改进点和经验", C_LIGHT),
    (7.5, "知识沉淀", "回流知识库\n更新规则表/Skill", C_GREEN),
    (10.5, "能力提升", "下次处置更准确\n更高效更有序", C_ORANGE),
]
for x, title, desc, color in loop:
    rect(s, x, 1.8, 2.2, 1.5, color)
    txt(s, x, 1.9, 2.2, 0.4, title, sz=16, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, 2.3, 2.2, 0.8, desc, sz=12, color=C_WHITE, align=PP_ALIGN.CENTER)
for x in [3.7, 6.7, 9.7]:
    txt(s, x, 2.2, 0.8, 0.5, "→", sz=28, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s, 12.7, 2.2, 0.5, 0.5, "↻", sz=28, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)

txt(s, 0.5, 3.6, 12, 0.4, "Review Agent 自动复盘内容", sz=18, color=C_DARK, bold=True)
review_items = [
    "① 时间线分析：各环节耗时、瓶颈识别、与标准时限对比",
    "② 决策回溯：定级是否合理、力量编成是否匹配、处置措施是否得当",
    "③ 信息流分析：哪些环节信息延迟/不一致、各组数据口径差异",
    "④ 改进建议：自动生成分部门整改建议，跟踪闭环",
    "⑤ 知识提取：新模式/新经验→更新知识库→规则表版本迭代→Skill升级",
]
mtxt(s, 0.5, 4.1, 12, 3, review_items, sz=14, color=C_GRAY)
pnum(s, 18)

# ==================== 第五章 ====================
chapter("五", "工程落地与安全审计", "188测试全绿 · 7 MCP工具 · Trace/Metrics · Guardrail")
pnum(s, 19)

# ==================== Slide 19: 完成度 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "当前工程完成度（188测试全绿）")

txt(s, 0.5, 1.1, 5, 0.4, "✅ 系统层（100%完成）", sz=18, color=C_GREEN, bold=True)
mtxt(s, 0.5, 1.6, 6, 3, [
    "✅ 12个功能模块全部开发完成",
    "✅ 规则引擎4张规则表v2026.2（含法规条文）",
    "✅ 三阶段状态机+完成标准校验",
    "✅ 多通道通知+强制反馈升级引擎",
    "✅ 5种事件类型包(INF/FOOD/ENV/POISON/UNK)",
    "✅ RBAC权限(11账号) · 75条回归测试全绿",
], sz=14, color=C_GREEN)

txt(s, 0.5, 4.0, 5, 0.4, "✅ Agent/Skill层（100%完成）", sz=18, color=C_GREEN, bold=True)
mtxt(s, 0.5, 4.5, 6, 3, [
    "✅ 7个职能Agent + 12个Skill v1.1.0",
    "✅ 双证据验证 + Guardrail + 证据链",
    "✅ 7个MCP工具 + Trace/Metrics可观测",
    "✅ 双模式运行 + AgentTeams适配层",
    "✅ SKILL.md×12 + WORKER.md×7 + AGENTS.md",
    "✅ 113条Agent测试全绿",
], sz=14, color=C_GREEN)

txt(s, 7, 1.1, 5, 0.4, "📊 测试覆盖", sz=18, color=C_ORANGE, bold=True)
mtxt(s, 7, 1.6, 5.5, 5, [
    "系统层 75条：",
    "  分级(D0) · 时限(D0) · 编成(D1) · 知识(D2)",
    "  边界情况 · 非法输入 · 规则未覆盖",
    "",
    "Agent层 113条：",
    "  12 Skill独立测试 · 7 Agent react测试",
    "  双证据验证 · 双模式切换 · MCP工具",
    "  Trace/Metrics可观测",
    "",
    "📌 合计 188条测试，全部通过",
    "📌 零新npm依赖",
], sz=14, color=C_GRAY)
pnum(s, 20)

# ==================== Slide 20: 安全 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "安全审计 + MCP工具 + 可观测体系")

sec = [
    ["机制", "实现"],
    ["RBAC权限", "11角色，单线决策（仅commander/deputy可启动/推进/终止）"],
    ["规则存证", "rule_results含ruleVersion+输入+输出+evidence证据链"],
    ["Guardrail", "12个Skill均含四类约束：确定性/业务/容错/安全"],
    ["证据链", "provenance/methodology/confidence/missing_evidence"],
    ["法规溯源", "legal_ref→国办函〔2006〕39号/国务院令第376号等"],
    ["通知审计", "全量通知+通道日志+升级记录，不可删除"],
    ["数据硬化", "原子写入+启动自检+每日快照（保留7份）"],
    ["高风险审批", "定级/启动/发布/终止均需人工确认（drill模式可代操作）"],
]
tbl(s, 0.3, 1.0, 6.5, 5, sec)

mcp = [
    ["MCP工具(7个)", "状态"],
    ["contacts(通讯录)", "✅ 已实现"],
    ["materials(物资)", "✅ 已实现"],
    ["vehicles(车辆)", "✅ 已实现"],
    ["medical-resources", "✅ 已实现"],
    ["knowledge(知识库)", "✅ 已实现"],
    ["notification-gateway", "✅ 已实现"],
    ["event-state(事件状态)", "✅ 已实现"],
]
tbl(s, 7.2, 1.0, 3.0, 4, mcp)

obs = [
    ["可观测", "状态"],
    ["Trace全链路", "✅ 已实现"],
    ["Metrics指标", "✅ 已实现"],
    ["Agent Sessions", "✅ 已实现"],
]
tbl(s, 10.5, 1.0, 2.5, 2.2, obs)

txt(s, 7.2, 5.2, 5.8, 0.8,
    "MCP标准协议: tools/list + tools/call\n"
    "可观测API: /api/agent/traces + /metrics",
    sz=12, color=C_MAIN)
pnum(s, 21)

# ==================== 第六章 ====================
chapter("六", "开放/开源计划", "MIT协议 · 可复用成果 · 第三方依赖")
pnum(s, 22)

# ==================== Slide 22: 开源 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "开源范围与可复用成果")
txt(s, 0.5, 0.9, 5, 0.4, "开源协议：MIT License", sz=16, color=C_GREEN, bold=True)

op = [
    ["开源范围", "内容"],
    ["多Agent协同框架", "AGENTS.md + WORKER.md×7 + 协同流程"],
    ["Skill工程化规范", "SKILL.md×12（Guardrail+证据链+references/）"],
    ["规则引擎", "4张规则表v2026.2（含法规条文引用）"],
    ["事件类型包", "5种事件类型配置（可扩展模板）"],
    ["测试套件", "188条测试用例（75系统+113 Agent）"],
    ["MCP工具", "7个标准MCP工具 + Trace/Metrics"],
    ["文档体系", "架构设计+接口契约+部署说明"],
]
tbl(s, 0.3, 1.4, 6.5, 4, op)

reuse = [
    ["可复用成果", "迁移场景"],
    ["Skill工程化规范", "任何需要可审计AI能力的场景"],
    ["双证据验证模式", "任何需要高可靠AI判断的场景"],
    ["Guardrail约束", "任何受监管行业的AI决策"],
    ["法规溯源体系", "规则表legal_ref→条文可追溯"],
    ["类型包插件化", "安全生产/自然灾害/城市应急"],
    ["双模式运行", "任何需要培训演练的场景"],
    ["强制反馈通知引擎", "运维告警/工单催办/审批催办"],
]
tbl(s, 7.2, 1.4, 5.8, 4, reuse)
pnum(s, 23)

# ==================== 第七章 ====================
chapter("七", "落地计划与进展", "当前进展 · 里程碑 · 风险控制")
pnum(s, 24)

# ==================== Slide 24: 里程碑 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "里程碑与风险控制")

mile = [
    ["阶段", "时间", "目标", "交付物"],
    ["初赛 ✅", "8.16", "方案设计+Agent/Skill层完成", "方案PPT ✅ · 188测试全绿 ✅"],
    ["复赛", "9.3", "AgentTeams桥接+运行验证", "可运行Demo+代码包+视频"],
    ["决赛", "9.22", "工程完善+现场展示", "路演PPT+现场Demo"],
    ["赛后", "Q4", "开源+实际场景验证", "GitHub仓库+社区"],
]
tbl(s, 0.5, 1.1, 12, 2.5, mile)

risk = [
    ["风险", "等级", "应对", "状态"],
    ["AgentTeams接入", "🟢 低", "standalone+bridge双模式已实现", "✅ 已解决"],
    ["MCP工具开发", "🟢 低", "7个MCP工具已实现，标准协议", "✅ 已解决"],
    ["可观测接入", "🟢 低", "Trace+Metrics已实现，API可查", "✅ 已解决"],
    ["闭源模型依赖", "🟢 低", "规则引擎兜底+LLM fallback降级", "✅ 已解决"],
    ["Skill法规准确性", "🟡 中", "references/存原文摘录+条文号", "✅ 已验证"],
]
tbl(s, 0.5, 4.0, 12, 2.8, risk)
pnum(s, 25)

# ==================== 第八章 ====================
chapter("八", "团队介绍", "成员背景 · 分工 · 成果")
pnum(s, 26)

# ==================== Slide 26: 团队 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, "团队介绍")

team = [
    ["成员", "角色", "背景", "核心技能"],
    ["【姓名】", "【队长/开发】", "【学校/公司】", "【全栈/架构】"],
    ["【姓名】", "【开发】", "【学校/公司】", "【前端/可视化】"],
    ["【姓名】", "【开发】", "【学校/公司】", "【后端/数据库】"],
]
tbl(s, 0.5, 1.2, 12, 2.5, team)

txt(s, 0.5, 4.0, 5, 0.4, "团队分工", sz=18, color=C_MAIN, bold=True)
mtxt(s, 0.5, 4.5, 12, 1.5, [
    "· 【成员1】：系统架构设计、规则引擎、后端核心",
    "· 【成员2】：前端开发、SVG图表库、交互设计",
    "· 【成员3】：数据层设计、测试套件、部署运维",
], sz=15, color=C_GRAY)

txt(s, 0.5, 5.8, 12, 0.4, "成果：12模块+7Agent+12Skill · 188条测试全绿 · 7MCP+Trace/Metrics · 双模式运行", sz=14, color=C_MAIN, bold=True)
pnum(s, 27)

# ==================== Slide 27: 尾页 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C_DARK)
rect(s, 0, 0, 13.333, 0.15, C_ORANGE)
txt(s, 1, 1.5, 11, 1, "卫盾Agent", sz=48, color=C_WHITE, bold=True)
txt(s, 1, 2.8, 11, 1.5,
    "7 Agent · 12 Skill · 双证据验证 · Guardrail约束\n法规溯源 · 188测试全绿 · 双模式运行",
    sz=22, color=C_SOFT)

hl = [
    ("🎯", "双证据验证", "LLM+规则引擎交叉校验"),
    ("🛡️", "Guardrail", "四类约束·证据链·可审计"),
    ("📋", "法规溯源", "legal_ref→具体条文可追溯"),
    ("🔄", "双模式", "辅助处置+模拟演练"),
]
for i, (icon, title, desc) in enumerate(hl):
    x = 1 + i * 3
    rect(s, x, 4.5, 2.5, 1.5, RGBColor(0x1E, 0x4A, 0x7A))
    txt(s, x, 4.55, 2.5, 0.4, icon, sz=26, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, x, 4.95, 2.5, 0.4, title, sz=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, 5.3, 2.5, 0.4, desc, sz=11, color=C_LBLUE, align=PP_ALIGN.CENTER)

txt(s, 1, 6.5, 11, 0.5, "让应急处置从「人盯人」走向「可协同·可治理·可复用」的Agent Infra", sz=16, color=C_LBLUE, align=PP_ALIGN.CENTER)
pnum(s, 28)

# ==================== 保存 ====================
out = r'E:\xiangmu\ceshi2\docs\competition\卫盾Agent_方案PPT.pptx'
prs.save(out)
print(f"PPT已生成: {out}")
print(f"共 {len(prs.slides)} 页")
